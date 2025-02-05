from pptx import Presentation
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
import pptx.shapes
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import comtypes.client
import os
import io


PPI_CONSTANT = 913080.788

# ----------- AUXILLARY FUNCTIONS -------------

def createLabel(context, label):
    new_label = tk.Label(context, text=label)
    new_label.pack(fill="y", padx=10, pady=10)
    return new_label

def createDropdown(context, optionsList):
    option_var = tk.StringVar()
    option_var.set(optionsList[0])
    option_menu = tk.OptionMenu(context, option_var, *optionsList)
    option_menu.pack(pady=5)

    return option_menu

# Function that formats a three long bulleted list into PowerPoint
def bulletedList(shape, list):
    shape.text_frame.clear()

    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = list[0]
    run.font.size = Pt(10)
    run.font.name = "Arial"
    run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    for item in list[1:]:
        p = shape.text_frame.add_paragraph()
        run = p.add_run()
        run.text = item
        run.font.size = Pt(10)
        run.font.name = "Arial"
        run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

def isPresentationOpen(ppt_app, fileName):
    prs = ppt_app.Presentations
    for pr in prs:
        print(pr.Name.lower())
        if pr.Name.lower() == fileName.lower():
            return True
    
    return False

# ---------------- APPLICATION -----------------------

# Class to define the GUI for the application
class TechProfileApp:
    def __init__(self, root):
        # Initialize, name, and format app window
        self.root = root
        self.root.title("Tech Profile Generator")
        self.root.geometry("1200x800")

        # Create title
        title = tk.Label(root, text="Tech Profile Generator", font=("Arial", 24))
        title.pack(pady=20)



        # Create frame for columns
        frame = tk.Frame(root)
        frame.pack(fill="both", padx=10, pady=10)

        # Create categories column
        categories_frame = tk.Frame(frame)
        categories_frame.pack(side="left", padx=10, pady=10)

        # Create dropdown months column
        options_frame = tk.Frame(frame)
        options_frame.pack(side="left", padx=10, pady=10)

        # Create text fields column
        text_fields_frame = tk.Frame(frame)
        text_fields_frame.pack(side="left", padx=10, pady=10)

        # Create initiatives and challenges column
        init_chall_frame = tk.Frame(frame)
        init_chall_frame.pack(side="left", padx=150, pady=10)



        # Populate category column
        self.categories = []
        self.category_text_fields = []
        self.init_text_fields = []
        self.chall_text_fields = []

        categoryTitle = createLabel(categories_frame, "Category")
        categoryTitle.config(font=("Arial", 12, "bold"), justify='center')

        self.categories.append(createLabel(categories_frame, "Workload 1"))
        self.categories.append(createLabel(categories_frame, "Workload 2"))
        self.categories.append(createLabel(categories_frame, "Workload 3"))
        self.categories.append(createLabel(categories_frame, "Virtualization"))
        self.categories.append(createLabel(categories_frame, "Compute"))
        self.categories.append(createLabel(categories_frame, "Networking"))
        self.categories.append(createLabel(categories_frame, "Storage"))
        self.categories.append(createLabel(categories_frame, "Replication"))
        self.categories.append(createLabel(categories_frame, "Backup"))

        # Populate dropdown column
        dropdownTitle = createLabel(options_frame, "Image Selection")
        dropdownTitle.config(font=("Arial", 12, "bold"), justify='center')

        workload = ["SQL Database", "Database (Generic)", "SAP", "File Server (Generic)", "Video Surveillance", "CAD (Generic)", "VDI (Generic)", "AutoCAD", "Kubernetes", "DHCP Server (Generic)", "Horizon/Omnissa", "Exchange", "Print Server (Generic)", "Office 365", "Oracle", "No Image"]
        virtualization = ["VMware", "HyperV", "No Image"]
        compute = ["PowerEdge 1U", "PowerEdge 2U", "VxRail 1U", "VxRail 2U", "No Image"]
        networking = ["S4128T", "S4148T", "S5224F", "S5248F", "S4128F", "S5232F", "S5296F", "Cisco BaseT 24 Port", "Cisco BaseT 48 Port", "Cisco FO Eth 24 Port", "Cisco FO Eth 48 Port", "HPE Aruba", "No Image"]
        storage = ["PowerStore", "Unity XT", "PowerVault", "HPE Storage", "NetApp Storage", "Pure Storage", "No Image"]
        replication = ["vSphere Replication", "Zerto", "No Image"]
        backup = ["Veeam", "Rubrik", "AWS", "Zerto", "Wasabi", "No Image"]

        self.workload_one_menu = createDropdown(options_frame, workload)
        self.workload_two_menu = createDropdown(options_frame, workload)
        self.workload_three_menu = createDropdown(options_frame, workload)
        self.virtualization_menu = createDropdown(options_frame, virtualization)
        self.compute_menu = createDropdown(options_frame, compute)
        self.networking_menu = createDropdown(options_frame, networking)
        self.storage_menu = createDropdown(options_frame, storage)
        self.replication_menu = createDropdown(options_frame, replication)
        self.backup_menu = createDropdown(options_frame, backup)

        # Populate text field column
        #TODO: Change text field size for multiple things
        textFieldTitle = createLabel(text_fields_frame, "Description")
        textFieldTitle.config(font=("Arial", 12, "bold"), justify='center')

        for i in range(9):
            text_field = tk.Text(text_fields_frame, height=1, width=30)
            text_field.pack(pady=11)
            self.category_text_fields.append(text_field)

        # Populate initiatives and challenges column
        init_title = createLabel(init_chall_frame, "Initiatives")
        init_title.config(font=("Arial", 12, "bold"), justify='center')
        for i in range(3):
            text_field = tk.Text(init_chall_frame, height=2, width=40)
            text_field.pack(pady=10)
            self.init_text_fields.append(text_field)

        chall_title = createLabel(init_chall_frame, "Challenges")
        chall_title.config(font=("Arial", 12, "bold"), justify='center')
        for i in range(3):
            text_field = tk.Text(init_chall_frame, height=2, width=40)
            text_field.pack(pady=10)
            self.chall_text_fields.append(text_field)
        


        # Make company name and file name section
        save_frame = tk.Frame(root)
        save_frame.pack(fill="both", expand=True, padx=10, pady=10)

        label_frame = tk.Frame(save_frame)
        label_frame.pack(side="left", fill="y", padx=10, pady=10)

        name_fields_frame = tk.Frame(save_frame)
        name_fields_frame.pack(side="left", fill="y", padx=10, pady=10)

        companyTitle = createLabel(label_frame, "Company Name")
        companyTitle.config(font=("Arial", 10, "bold"))
        fileTitle = createLabel(label_frame, "File Name")
        fileTitle.config(font=("Arial", 10, "bold"))

        self.company_name = tk.Text(name_fields_frame, height=1, width=30)
        self.company_name.pack(fill="y", padx=10, pady=12)
        
        self.file_name = tk.Text(name_fields_frame, height=1, width=30)
        self.file_name.pack(fill="y", padx=10, pady=12)



        # Make generate button
        generate_button = tk.Button(root, text="Generate", command=self.generate_tech_profile)
        generate_button.pack(pady=20)


    # Event handler that generates tech profile once button is pressed
    def generate_tech_profile(self):
        ppt_app = comtypes.client.CreateObject("PowerPoint.Application")

        if isPresentationOpen(ppt_app, "template_v1.pptx"):
            messagebox.showinfo("Error", "The template this program uses is currently open. Please close the template and try again.")
            return

        # Load tech profile template
        tp = Presentation("template_v1.pptx")
        

        # ------------------ STEP 1 ------------------------
        # Parse text from text fields and place in variables


        # Parse file name and error check for if the file is open
        fileName = self.file_name.get("1.0", tk.END).strip('\n')
        print(fileName)

        if isPresentationOpen(ppt_app, fileName + ".pptx"):
            messagebox.showinfo("Error", "The file you are trying to create/update is currently open. Please close and try again.")
            return
        
        if fileName == "":
            messagebox.showinfo("Error", "Please enter a name for your tech profile")
            return

        # Parse company name
        companyName = self.company_name.get("1.0", tk.END).strip('\n')
        print(companyName)

        # Parse initiatives
        initiatives = []
        for i in range(3):
            init = self.init_text_fields[i].get('1.0', tk.END).strip('\n')
            print(init)
            initiatives.append(init)

        # Parse challenges
        challenges = []
        for i in range(3):
            chall = self.chall_text_fields[i].get('1.0', tk.END).strip('\n')
            print(chall)
            challenges.append(chall)

        # Parse workloads
        workloads = []
        for i in range(3):
            work = self.category_text_fields[i].get('1.0', tk.END).strip('\n')
            print(work)
            workloads.append(work)

        # Parse virtualization, compute, network, storage, replication, and backup
        v = self.category_text_fields[3].get('1.0', tk.END).strip('\n')
        c = self.category_text_fields[4].get('1.0', tk.END).strip('\n')
        n = self.category_text_fields[5].get('1.0', tk.END).strip('\n')
        s = self.category_text_fields[6].get('1.0', tk.END).strip('\n')
        r = self.category_text_fields[7].get('1.0', tk.END).strip('\n')
        b = self.category_text_fields[8].get('1.0', tk.END).strip('\n')
        print(v,c,n,s,r,b)


        # ------------------- STEP 2 -----------------------
        # Load tech profile slide, read for each textbox, and place the appropriate user input into the field


        slide = tp.slides[0]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "Single Site" in run.text:
                        run.text = companyName
                    elif "Edit I" in run.text:
                        bulletedList(shape, initiatives)
                    elif "Ch Edit" in run.text:
                        bulletedList(shape, challenges)
                    elif "Workload 1" in run.text:
                        run.text = workloads[0]
                    elif "Workload 2" in run.text:
                        run.text = workloads[1]
                    elif "Workload 3" in run.text:
                        run.text = workloads[2]
                    elif "Edit V" in run.text:
                        run.text = v
                    elif "Edit C" in run.text:
                        run.text = c
                    elif "Edit N" in run.text:
                        run.text = n
                    elif "Edit S" in run.text:
                        run.text = s
                    elif "Edit R" in run.text:
                        run.text = r
                    elif "Edit B" in run.text:
                        run.text = b


        # ----------------- STEP 3 ---------------------
        # Parse dropdown menu selections


        w1Selection = self.workload_one_menu.cget("text")
        w2Selection = self.workload_two_menu.cget("text")
        w3Selection = self.workload_three_menu.cget("text")
        vSelection = self.virtualization_menu.cget("text")
        cSelection = self.compute_menu.cget("text")
        nSelection = self.networking_menu.cget("text")
        sSelection = self.storage_menu.cget("text")
        rSelection = self.replication_menu.cget("text")
        bSelection = self.backup_menu.cget("text")

        # print(w1Selection)
        # print(w2Selection)
        # print(w3Selection)
        # print(vSelection)
        # print(cSelection)
        # print(nSelection)
        # print(sSelection)
        # print(rSelection)
        # print(bSelection)


        # ----------------- STEP 4 ------------------------
        # Copy image from template slide based on saved dropdown keyword, paste into new slide, and apply formatting
        #
        # 1. Find the image in respective slide (e.g. servers in Compute slide) through matching the saved dropdown keyword to the name of the shape of the image
        # 2. Add image to first slide of template slide deck
        # 3. Change position values of image

        # Workloads
        workloadSlide = tp.slides[1]
        workloadImage1 = workloadSlide.shapes[0]
        workloadImage2 = workloadSlide.shapes[0]
        workloadImage3 = workloadSlide.shapes[0]

        for shape in workloadSlide.shapes:
            if shape.name == w1Selection:
                workloadImage1 = shape
            if shape.name == w2Selection:
                workloadImage2 = shape
            if shape.name == w3Selection:
                workloadImage3 = shape
        
        # print(workloadImage1.name)

        if workloadImage1 != workloadSlide.shapes[0]:
            imageStream = io.BytesIO(workloadImage1.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(1.89), top=Inches(1.37), width=workloadImage1.width, height=workloadImage1.height)

        if workloadImage2 != workloadSlide.shapes[0]:
            imageStream = io.BytesIO(workloadImage2.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(4.94), top=Inches(1.37), width=workloadImage2.width, height=workloadImage2.height)

        if workloadImage3 != workloadSlide.shapes[0]:
            imageStream = io.BytesIO(workloadImage3.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(8.02), top=Inches(1.37), width=workloadImage3.width, height=workloadImage3.height)

        # Virtualization
        virtualizationSlide = tp.slides[2]
        virtualizationImage = virtualizationSlide.shapes[0]

        for shape in virtualizationSlide.shapes:
            if shape.name == vSelection:
                virtualizationImage = shape
        
        # print(virtualizationImage.name)
        
        if virtualizationImage != virtualizationSlide.shapes[0]:
            imageStream = io.BytesIO(virtualizationImage.image.blob)

            if "VM" in virtualizationImage.name:
                slide.shapes.add_picture(imageStream, left=Inches(3.18), top=Inches(2.29), width=virtualizationImage.width, height=virtualizationImage.height)
            else:
                slide.shapes.add_picture(imageStream, left=Inches(2.86), top=Inches(2.4), width=virtualizationImage.width, height=virtualizationImage.height)


        # Compute
        computeSlide = tp.slides[3]
        computeImage = computeSlide.shapes[0]

        for shape in computeSlide.shapes:
            if shape.name == cSelection:
                computeImage = shape
        
        # print(computeImage.name)
        
        if computeImage != computeSlide.shapes[0]:
            imageStream = io.BytesIO(computeImage.image.blob)

            if "2U" in computeImage.name:
                slide.shapes.add_picture(imageStream, left=Inches(1.98), top=Inches(3.38), width=computeImage.width, height=computeImage.height)
            else:
                slide.shapes.add_picture(imageStream, left=Inches(1.98), top=Inches(3.52), width=computeImage.width, height=computeImage.height)


        # Networking
        networkingSlide = tp.slides[4]
        networkingImage = networkingSlide.shapes[0]

        for shape in networkingSlide.shapes:
            if shape.name == nSelection:
                networkingImage = shape
        
        # print(networkingImage.name)
        
        if networkingImage != networkingSlide.shapes[0]:
            imageStream = io.BytesIO(networkingImage.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(2), top=Inches(4.48), width=networkingImage.width, height=networkingImage.height)


        # Storage
        storageSlide = tp.slides[5]
        storageImage = storageSlide.shapes[0]

        for shape in storageSlide.shapes:
            if shape.name == sSelection:
                storageImage = shape
        
        # print(storageImage.name)
        
        if storageImage != storageSlide.shapes[0]:
            imageStream = io.BytesIO(storageImage.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(2.02), top=Inches(5.41), width=storageImage.width, height=storageImage.height)


        # Replication
        replicationSlide = tp.slides[6]
        replicationImage = replicationSlide.shapes[0]

        for shape in replicationSlide.shapes:
            if shape.name == rSelection:
                replicationImage = shape
        
        # print(replicationImage.name)
        
        if replicationImage != replicationSlide.shapes[0]:
            imageStream = io.BytesIO(replicationImage.image.blob)

            slide.shapes.add_picture(imageStream, left=Inches(1.37), top=Inches(6.5), width=replicationImage.width, height=replicationImage.height)


        # Backup
        backupSlide = tp.slides[7]
        backupImage = backupSlide.shapes[0]

        for shape in backupSlide.shapes:
            if shape.name == bSelection:
                backupImage = shape
        
        # print(backupImage.name)
        
        if backupImage != backupSlide.shapes[0]:
            imageStream = io.BytesIO(backupImage.image.blob)

            print(backupImage.width)
            print(backupImage.height)

            slide.shapes.add_picture(imageStream, left=Inches(4.28), top=Inches(6.5), width=backupImage.width, height=backupImage.height)


        # ----------------- STEP 4 ----------------------
        # Save new PPT, reopen and delete slides

        # Save tech profile
        home_directory = os.path.expanduser("~")

        target_folder = os.path.join(home_directory, "Downloads")
        # target_folder = filedialog.askdirectory(title="Select Folder to Save Presentation")

        # if not target_folder:
        #     messagebox.showinfo("Error", "No folder selected. Please try again.")
        #     return

        new_prs_name = fileName + ".pptx"
        save_path = os.path.join(target_folder, new_prs_name)

        if os.path.exists(save_path):
            response = messagebox.askyesno("File already exists", "The file you are trying to create already exists. Do you want to overwrite it?")

            if response:
                tp.save(save_path)
            else:
                return            
        else:
            tp.save(save_path)
        
        # Delete slides
        # if getattr(sys, 'frozen', False):
        #     save_path = os.path.join(os.path.dirname(sys.executable), fileName + ".pptx")
        # else:
        #     save_path = os.path.join(os.path.dirname(__file__), fileName + ".pptx")

        if os.path.exists(save_path):
            try:
                prs = ppt_app.Presentations.Open(save_path, WithWindow=False)

                totalSlides = prs.Slides.Count

                slidesToDelete = []
                for i in range(2, totalSlides + 1):
                    slidesToDelete.append(i)
                
                for i in reversed(slidesToDelete):
                    prs.Slides(i).Delete()
                
                prs.Save()

                prs.Close()
            except Exception as e:
                print(f"Error: {e}")
        else:
            messagebox.showerror("Error", "Failed to save the presentation. Cannot open it.")
        
        # ----------------- STEP 5 ------------------
        # Post confirmation dialog

        messagebox.showinfo("Created!", "Your file was created successfully!")

    def run(self):
        self.root.mainloop()


""" Steps to load image into tech profile:
        1. Design GUI with the three columns, one having drop down options
        2. Match drop down options to pulling correct image from template slides
        3. Write function that copies image, pastes into template slide, then change size and position to preset value based on category
        4. Adjust previous code to work with new GUI
        5. Adjust code to remove SOTEC slides from template before saving
"""

# TODO: Add a README for instructions
# TODO: Make a template for autofill
# TODO: Implement function that can auto fill text fields based on a text file input


if __name__ == "__main__":
    root = tk.Tk()
    app = TechProfileApp(root)
    app.run()

