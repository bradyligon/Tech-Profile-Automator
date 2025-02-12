# IO Libraries
import comtypes.client
import os
import io

# PowerPoint Libraries
from pptx import Presentation
import pptx.shapes
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

# Custom Libraries
from component_data import *



# -------------------- AUXILLARY FUNCTIONS ---------------------------

def update_paragraph(paragraph, run_text, isBold):
    # Function that updates paragraph with new text
    paragraph.clear()

    run = paragraph.add_run()
    run.text = run_text

    font = run.font
    font.name = "Arial"
    font.size = Pt(8)
    font.bold = isBold
    font.color.rgb = RGBColor(0, 0, 0)

    return run

def load_image_data(slide, imageNames):
    # Function that saves the size and position data of image in a dictionary, given a slide from a PowerPoint and a list of names to identify specific images to save
    imageData = {}

    for shape in slide.shapes:
        if shape.name in imageNames:
            imageData[shape.name] = {
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            }

    return imageData



# -------------------- MAIN FUNCTION -------------------------

def generate_tech_profile():
    # Generates the tech profile. Below are parameters to add in later:
    # initTxt, challTxt, workloadTxt, virtualizationInput, computeInput, networkInput, managementInput, storageInput, replicationInput, backupInput
    numSites = 1
    networkInput = Component("PowerSwitch S5224F", 2)
    managementInput = Component("PowerSwitch S4128T", 1)
    computeInput = []
    computeInput.append(Compute(Component("PowerEdge R660", 3), Component("Intel Gold 6526Y, 2.8 GHz, 16C/32T", 1), "512 GiB"))

    # Load template presentation
    tp = Presentation("template_v2.pptx")


    # ------------------- LOAD NECESSARY DATA ---------------------
    # -------------------------------------------------------------

    singleImageNames = ["Switch1", "Switch2", "ManagementSwitch", "Server1", "Server2", "Server3", "Server4", "Server5", "Server6", "Server7", "Server8", "Server9", "Server10", "Server11", "Server12", "Storage", "Backup", "Virtualization"]
    doubleImageNames = ["Switch1_1", "Switch2_1", "Switch1_2", "Switch2_2", "ManagementSwitch1", "ManagementSwitch2", "Server1_1", "Server2_1", "Server3_1", "Server4_1", "Server5_1", "Server6_1", "Server1_2", "Server2_2", "Server3_2", "Server4_2", "Server5_2", "Server6_2", "Storage1", "Storage2", "Backup1", "Backup2", "Virtualization"]

    # Load images' left, top, width, and height data from the slide based on if this is a single or double site setup
    slide = tp.slides[numSites - 1]

    if numSites == 1:
        imageData = load_image_data(slide, singleImageNames)

        # Print the results (for debugging)
        # for name, dims in imageData.items():
        #     print(f"Shape: {name}")
        #     print(f"  Left: {dims['left']}")
        #     print(f"  Top: {dims['top']}")
        #     print(f"  Width: {dims['width']}")
        #     print(f"  Height: {dims['height']}")

    else:
        imageData = load_image_data(slide, doubleImageNames)

        # Print the results (for debugging)
        # for name, dims in imageData.items():
        #     print(f"Shape: {name}")
        #     print(f"  Left: {dims['left']}")
        #     print(f"  Top: {dims['top']}")
        #     print(f"  Width: {dims['width']}")
        #     print(f"  Height: {dims['height']}")

    # Load switch data
    switches = gen_switch_data()


    # --------------------- EDIT TEXT FIELDS ----------------------
    # -------------------------------------------------------------

    for shape in slide.shapes:
        if shape.name == "Network_Info":
            # Edit text field for networking
            paragraphs = shape.text_frame.paragraphs

            for name in switches:
                if networkInput.model == name:
                    # ToR Model Description
                    switchModelTxt = "(" + str(networkInput.qty) + ") " + networkInput.model
                    update_paragraph(paragraphs[1], switchModelTxt, True)

                    # Speed Description
                    switchSpeedTxt = switches[name]['speed']
                    update_paragraph(paragraphs[2], switchSpeedTxt, False)

                    # Num Ports Description
                    switchPortsTxt = switches[name]['ports']
                    update_paragraph(paragraphs[3], switchPortsTxt, False)
                if managementInput.model == name:
                    # Management Switch Description
                    managementTxt = "(" + str(managementInput.qty) + ") " + managementInput.model
                    update_paragraph(paragraphs[5], managementTxt, True)

        if shape.name == "Compute_Info":
            # Edit text field for compute
            paragraphs = shape.text_frame.paragraphs

            row = 0
            for compute in computeInput:
                # Server Model Description
                serverTxt = "(" + str(compute.server.qty) + ") " + compute.server.model  
                update_paragraph(paragraphs[row], serverTxt, True)
                row = row + 1

                # CPU Description
                cpuTxt = "(" + str(compute.cpu.qty) + ") " + compute.cpu.model
                update_paragraph(paragraphs[row], cpuTxt, False)
                row = row + 1

                # Memory Description
                memoryTxt = compute.memory
                update_paragraph(paragraphs[row], memoryTxt, False)
                row = row + 1

            for i in range(row, 9):
                paragraphs[i].clear()

        
            



    # --------------------- COPY AND PAST IMAGES ----------------------
    # ----------------------------------------------------------------- 



    # Save tech profile
    fileName = "test"

    home_directory = os.path.expanduser("~")

    target_folder = os.path.join(home_directory, "Downloads")

    new_prs_name = fileName + ".pptx"
    save_path = os.path.join(target_folder, new_prs_name)

    tp.save(save_path)

    return

if __name__ == "__main__":
    generate_tech_profile()
    print("Done!")